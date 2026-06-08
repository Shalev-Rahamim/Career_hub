import uuid
import io
import os
import shutil
import tempfile
from fastapi import APIRouter, Depends, HTTPException, Header, UploadFile, File, Form
from sqlalchemy.orm import Session
from typing import List, Optional
import chromadb
import fitz  # PyMuPDF
import docx
import pptx
import zipfile

# Google GenAI SDK imports
from google import genai
from google.genai import types
from pydantic import BaseModel, Field

from app.database import get_db
from app.models import Interview, InterviewEvaluation, User
from app.schemas import (
    GenerateQuestionsResponse,
    QuestionItem,
    EvaluateInterviewRequest,
    EvaluateInterviewResponse,
    QuestionEvaluationItem
)
from app.config import settings

router = APIRouter(prefix="/api/interview", tags=["Interview Simulator"])

# --- Pydantic Schemas for LLM Structured Output ---
class LLMQuestionItem(BaseModel):
    question_text: str = Field(description="The tailored technical project defense question in the requested language")
    category: str = Field(description="Category of the question (e.g., Architecture, Security, Scaling) in the requested language")

class LLMQuestionsList(BaseModel):
    questions: List[LLMQuestionItem] = Field(description="List of 5 to 10 synthesized questions")

class LLMAnswerEvaluation(BaseModel):
    question_id: str = Field(description="ID of the question")
    score: int = Field(description="Score for this specific answer (0-100)", ge=0, le=100)
    rationale: str = Field(description="Constructive critique detailing what was good and what was missing in the requested language")
    model_answer: str = Field(description="An exemplar response representing a perfect answer for this question in the requested language")
    improved_phrasing: str = Field(description="The candidate's answer re-phrased to sound highly professional and technical in the requested language")

class LLMEvaluationResult(BaseModel):
    overall_score: int = Field(description="Average score across all questions (1-100)")
    general_feedback: str = Field(description="High-level feedback of candidate's strength and overall performance in the requested language")
    evaluations: List[LLMAnswerEvaluation] = Field(description="Detailed grading matrix for each answer")


# Initialize local Chroma client
def get_chroma_collection():
    try:
        client = chromadb.PersistentClient(path=settings.CHROMA_DB_DIR)
        collection = client.get_or_create_collection(
            name="system_evaluation_questions"
        )
        return collection
    except Exception as e:
        client = chromadb.EphemeralClient()
        collection = client.get_or_create_collection(
            name="system_evaluation_questions"
        )
        return collection

# Seeding list (Default Bank of 16 Technical/Architectural Defense Questions)
DEFAULT_QUESTION_BANK = [
    {
        "text": "In a stateful client-side application, race conditions can occur if multiple concurrent actions update state. How does your code handle state locking or transitions?",
        "category": "State Management"
    },
    {
        "text": "Your local SQLite database is a single point of failure (SPOF). How would you transition this system to support multi-region replication and failover?",
        "category": "Database Scaling"
    },
    {
        "text": "The codebase connects to external AI APIs synchronously or blockingly. How does your backend handle asynchronous connection limits and request timeouts?",
        "category": "Asynchronous Design"
    },
    {
        "text": "If the vector database is temporarily unavailable during question generation, how does the system degrade gracefully?",
        "category": "Error Fallbacks"
    },
    {
        "text": "If a component crashes during a database transaction, how do you prevent data inconsistency and ensure ACID compliance?",
        "category": "Data Consistency"
    },
    {
        "text": "How do you protect the backend from memory leaks during large file parses or high-throughput JSON processing?",
        "category": "Resource Management"
    },
    {
        "text": "How does your architecture prevent SQL injection and dynamic query vulnerability in search features?",
        "category": "Security"
    },
    {
        "text": "How does the system defend against Prompt Injection or jailbreaks targeting LLM endpoints?",
        "category": "AI Security"
    },
    {
        "text": "How do you handle API key rotation and credential containment in distributed microservices?",
        "category": "Infrastructure Security"
    },
    {
        "text": "Under 10000 concurrent requests, how do you prevent database thread starvation and connection pool depletion?",
        "category": "High Concurrency"
    },
    {
        "text": "Explain your strategy for caching frequently-queried vectors to reduce database load and query latency.",
        "category": "Optimization"
    },
    {
        "text": "What design patterns did you employ to decouple core features and enable independent deployment cycles?",
        "category": "Architectural Patterns"
    },
    {
        "text": "If a user cancels a long-running request mid-process, how do you ensure server resources are immediately released?",
        "category": "Resource Lifecycle"
    },
    {
        "text": "How do you track and audit user actions and system logs in compliance with data privacy regulations?",
        "category": "Logging & Auditing"
    },
    {
        "text": "In case of total model api failure, how does the system fallback to local mock generators or cached responses?",
        "category": "Resilience"
    },
    {
        "text": "Explain the strategy used to version-control database schema updates alongside code modifications without downtime.",
        "category": "Database Migrations"
    }
]

# Helper to populate Chroma DB
def seed_questions_into_chroma():
    collection = get_chroma_collection()
    if collection.count() == 0:
        try:
            # We can use LangChain Embeddings or native GenAI for seeding
            from langchain_google_genai import GoogleGenerativeAIEmbeddings
            embeddings = GoogleGenerativeAIEmbeddings(
                model="models/text-embedding-004",
                google_api_key=settings.GOOGLE_API_KEY
            )
            ids = [f"bank-{i}" for i in range(len(DEFAULT_QUESTION_BANK))]
            documents = [q["text"] for q in DEFAULT_QUESTION_BANK]
            metadatas = [{"category": q["category"]} for q in DEFAULT_QUESTION_BANK]
            embedded_docs = embeddings.embed_documents(documents)
            collection.add(
                ids=ids,
                embeddings=embedded_docs,
                documents=documents,
                metadatas=metadatas
            )
        except Exception as e:
            # Fallback mock setup
            ids = [f"bank-{i}" for i in range(len(DEFAULT_QUESTION_BANK))]
            documents = [q["text"] for q in DEFAULT_QUESTION_BANK]
            metadatas = [{"category": q["category"]} for q in DEFAULT_QUESTION_BANK]
            collection.add(
                ids=ids,
                documents=documents,
                metadatas=metadatas
            )


# Local file parser helper for Chroma DB querying
async def extract_text_from_upload_file(upload_file: UploadFile) -> str:
    filename = upload_file.filename.lower()
    file_bytes = await upload_file.read()
    # Seek back to 0 so other streams can read it
    await upload_file.seek(0)
    
    try:
        if filename.endswith(".pdf"):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            return "".join([page.get_text() for page in doc])
        elif filename.endswith((".docx", ".doc")):
            doc_file = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc_file.paragraphs]
            for table in doc_file.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paragraphs.append(cell.text)
            return "\n".join(paragraphs)
        elif filename.endswith((".pptx", ".ppt")):
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        else:
            return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return ""


# Helper to stream a file to Google GenAI File API
async def upload_to_google_ai(client: genai.Client, upload_file: UploadFile) -> any:
    # Write to a temporary file locally so google-genai SDK can stream it
    suffix = os.path.splitext(upload_file.filename)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
        content = await upload_file.read()
        await upload_file.seek(0)
        temp_file.write(content)
        temp_file_path = temp_file.name

    try:
        # Upload
        uploaded_file = client.files.upload(file=temp_file_path)
        return uploaded_file
    finally:
        # Clean up local temp file
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)


async def process_and_upload_file(client: genai.Client, upload_file: UploadFile) -> List[str]:
    filename = upload_file.filename.lower()
    
    # 1. If it's a PDF, TXT, or direct image, we can upload it directly.
    supported_direct_extensions = (".pdf", ".txt", ".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif", ".gif")
    if filename.endswith(supported_direct_extensions):
        uploaded = await upload_to_google_ai(client, upload_file)
        return [uploaded.name]
        
    # 2. For unsupported document types (Word/PowerPoint), extract text locally, upload as a .txt file,
    # and extract embedded media files to upload them individually.
    resource_names = []
    
    # Extract text content
    extracted_text = await extract_text_from_upload_file(upload_file)
    if not extracted_text.strip():
        extracted_text = f"Empty text or unparseable document: {upload_file.filename}"
        
    # Write the extracted text into a temporary TXT file and upload it
    with tempfile.NamedTemporaryFile(delete=False, suffix=".txt", mode="w", encoding="utf-8") as temp_txt:
        temp_txt.write(extracted_text)
        temp_txt_path = temp_txt.name
        
    try:
        uploaded_txt = client.files.upload(file=temp_txt_path)
        resource_names.append(uploaded_txt.name)
    finally:
        if os.path.exists(temp_txt_path):
            os.remove(temp_txt_path)
            
    # Now try to open the file as a zip and extract media files
    try:
        file_bytes = await upload_file.read()
        await upload_file.seek(0)
        
        # Open bytes as a zip archive
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for zip_info in z.infolist():
                name = zip_info.filename.lower()
                is_image = any(name.endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"))
                is_media_path = "media/" in name or "word/media/" in name or "ppt/media/" in name
                
                if is_image and (is_media_path or zip_info.file_size > 0):
                    img_data = z.read(zip_info.filename)
                    suffix = os.path.splitext(name)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_img:
                        temp_img.write(img_data)
                        temp_img_path = temp_img.name
                        
                    try:
                        uploaded_img = client.files.upload(file=temp_img_path)
                        resource_names.append(uploaded_img.name)
                    except Exception as e:
                        print(f"Skipping image {name} upload failure: {e}")
                    finally:
                        if os.path.exists(temp_img_path):
                            os.remove(temp_img_path)
    except Exception as e:
        print(f"File {upload_file.filename} is not a valid zip archive or zip extraction failed: {e}")
        
    return resource_names


@router.post("/generate-questions", response_model=GenerateQuestionsResponse)
async def generate_defense_questions(
    assignment_file: UploadFile = File(...),
    solution_file: UploadFile = File(...),
    difficulty_level: str = Form("medium"),
    num_questions: int = Form(5),
    language: str = Form("hebrew"),
    x_user_id: Optional[str] = Header(None, alias="X-User-ID"),
    db: Session = Depends(get_db)
):
    # Ensure user exists
    user_uuid = x_user_id or str(uuid.uuid4())
    user = db.query(User).filter(User.user_id == user_uuid).first()
    if not user:
        user = User(user_id=user_uuid)
        db.add(user)
        db.commit()
        db.refresh(user)

    # 1. Parse a small snippet locally to run Chroma DB vector similarity search
    retrieved_questions = []
    try:
        assignment_text_snippet = await extract_text_from_upload_file(assignment_file)
        
        collection = get_chroma_collection()
        has_key = settings.GOOGLE_API_KEY is not None
        
        if has_key and assignment_text_snippet.strip():
            from langchain_google_genai import GoogleGenerativeAIEmbeddings
            embeddings = GoogleGenerativeAIEmbeddings(
                model="models/text-embedding-004",
                google_api_key=settings.GOOGLE_API_KEY
            )
            query_vector = embeddings.embed_query(assignment_text_snippet[:2000])
            results = collection.query(
                query_embeddings=[query_vector],
                n_results=5
            )
        else:
            results = collection.query(
                query_texts=[assignment_text_snippet[:1000]] if assignment_text_snippet else ["General architecture"],
                n_results=5
            )

        if results and "documents" in results and results["documents"]:
            docs = results["documents"][0]
            metas = results["metadatas"][0] if "metadatas" in results else []
            for i in range(len(docs)):
                category = metas[i]["category"] if i < len(metas) else "General"
                retrieved_questions.append({
                    "question_text": docs[i],
                    "category": category
                })
    except Exception as e:
        for q in DEFAULT_QUESTION_BANK[:5]:
            retrieved_questions.append({"question_text": q["text"], "category": q["category"]})

    if not retrieved_questions:
        for q in DEFAULT_QUESTION_BANK[:5]:
            retrieved_questions.append({"question_text": q["text"], "category": q["category"]})

    # 2. Upload files to Google AI platform and generate structured questions natively
    try:
        client = genai.Client(api_key=settings.GOOGLE_API_KEY)
        
        # Upload dynamically
        assignment_uris = await process_and_upload_file(client, assignment_file)
        solution_uris = await process_and_upload_file(client, solution_file)
        
        assignment_files = []
        for uri in assignment_uris:
            try:
                assignment_files.append(client.files.get(name=uri))
            except Exception as e:
                print(f"Error getting assignment file {uri}: {e}")
                
        solution_files = []
        for uri in solution_uris:
            try:
                solution_files.append(client.files.get(name=uri))
            except Exception as e:
                print(f"Error getting solution file {uri}: {e}")
        
        # Call native Gemini model
        lang_str = "Hebrew" if language.lower() == "hebrew" else "English"
        system_prompt = f"You are a senior tech lead reviewing a candidate's technical home assignment.\n" \
                        f"You are given the uploaded assignment guidelines file(s) and the candidate's solution file(s).\n" \
                        f"Synthesize exactly {num_questions} customized project defense questions.\n" \
                        f"Focus the difficulty of the questions on a level appropriate for: {difficulty_level.upper()}.\n" \
                        f"- EASY: Focus on basic code clarity, syntax, simple logic, and basic local error handling.\n" \
                        f"- MEDIUM: Focus on standard design patterns, API separation, database usage, clean code, and standard testability.\n" \
                        f"- HARD: Focus on deep architectural patterns, high concurrency, security under pressure, scaling bottlenecks, memory leakage, and complex performance trade-offs.\n" \
                        f"For each question:\n" \
                        f"- Blend the retrieved conceptual query with the candidate's actual submitted solution/instructions.\n" \
                        f"- Write the questions clearly in {lang_str}."

        user_prompt = f"Retrieved Base System Questions to integrate or draw themes from:\n{str(retrieved_questions)}\n\n" \
                      f"Synthesize exactly {num_questions} questions in {lang_str}."

        # Model call
        contents_payload = []
        contents_payload.extend(assignment_files)
        contents_payload.extend(solution_files)
        contents_payload.append(system_prompt)
        contents_payload.append(user_prompt)

        response = client.models.generate_content(
            model=settings.LLM_MODEL or "gemini-2.0-flash",
            contents=contents_payload,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=LLMQuestionsList,
                temperature=0.4
            )
        )

        import json
        llm_result = LLMQuestionsList.model_validate_json(response.text)
        
        questions_list = []
        for idx, q in enumerate(llm_result.questions[:num_questions]):
            questions_list.append(QuestionItem(
                question_id=f"q-{idx + 1}",
                question_text=q.question_text,
                category=q.category
            ))

        # Save Interview session to database (comma-separated URIs)
        db_interview = Interview(
            user_id=user.user_id,
            assignment_file_uri=",".join(assignment_uris),
            solution_file_uri=",".join(solution_uris),
            difficulty_level=difficulty_level,
            num_questions=num_questions,
            language=language,
            questions_json=[q.model_dump() for q in questions_list]
        )
        db.add(db_interview)
        db.commit()
        db.refresh(db_interview)

        return GenerateQuestionsResponse(
            interview_id=db_interview.interview_id,
            questions=questions_list
        )

    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to generate defense questions: {str(e)}")


@router.post("/evaluate", response_model=EvaluateInterviewResponse)
async def evaluate_defense_answers(
    payload: EvaluateInterviewRequest,
    db: Session = Depends(get_db)
):
    # Fetch interview row from DB
    interview = db.query(Interview).filter(Interview.interview_id == payload.interview_id).first()
    if not interview:
        raise HTTPException(status_code=404, detail="Interview session not found")

    answers_mapping = {ans.question_id: ans.answer_text for ans in payload.answers}
    questions_list = interview.questions_json or []

    conversation_str = ""
    for q in questions_list:
        q_id = q.get("question_id")
        q_text = q.get("question_text")
        ans_text = answers_mapping.get(q_id, "No answer provided.")
        conversation_str += f"Question ID: {q_id}\nQuestion: {q_text}\nCandidate's Answer: {ans_text}\n\n"

    try:
        client = genai.Client(api_key=settings.GOOGLE_API_KEY)
        
        # Retrieve original file references from Google AI Files API dynamically
        lang = getattr(interview, 'language', 'hebrew') or 'hebrew'
        lang_str = "Hebrew" if lang.lower() == "hebrew" else "English"
        
        assignment_uris = (interview.assignment_file_uri or "").split(",")
        solution_uris = (interview.solution_file_uri or "").split(",")
        
        assignment_files = []
        for uri in assignment_uris:
            uri = uri.strip()
            if uri:
                try:
                    assignment_files.append(client.files.get(name=uri))
                except Exception as e:
                    print(f"Error getting assignment file {uri}: {e}")
                    
        solution_files = []
        for uri in solution_uris:
            uri = uri.strip()
            if uri:
                try:
                    solution_files.append(client.files.get(name=uri))
                except Exception as e:
                    print(f"Error getting solution file {uri}: {e}")

        system_prompt = f"You are a senior tech lead evaluating a candidate's technical defense of their home assignment.\n" \
                        f"Compare the candidate's answers against the context of the attached assignment instructions file and solution file.\n" \
                        f"The chosen difficulty is: {getattr(interview, 'difficulty_level', 'medium').upper()}.\n" \
                        f"Evaluate the answers with strictness matching this difficulty level:\n" \
                        f"- EASY: Evaluate if the answers explain basic programming concepts, syntax, and simple logic correctness.\n" \
                        f"- MEDIUM: Evaluate if the answers demonstrate standard design principles, appropriate API splits, and clean code.\n" \
                        f"- HARD: Expect senior-level reasoning showing deep architectural patterns, concurrency knowledge, performance benchmarking, and complex trade-off analysis.\n" \
                        f"For each answer:\n" \
                        f"- Score it between 0 and 100.\n" \
                        f"- Provide a detailed rationale/feedback strictly in {lang_str}.\n" \
                        f"- Provide an ideal model answer in {lang_str}.\n" \
                        f"- Provide an improved, highly technical phrasing in {lang_str} that the candidate could use to sound more senior.\n\n" \
                        f"Also provide an overall score (1-100) and general high-level feedback strictly in {lang_str}."

        user_prompt = f"Candidate's Answers:\n{conversation_str}"

        # Generate evaluation
        contents_payload = []
        contents_payload.extend(assignment_files)
        contents_payload.extend(solution_files)
        contents_payload.append(system_prompt)
        contents_payload.append(user_prompt)

        response = client.models.generate_content(
            model=settings.LLM_MODEL or "gemini-2.0-flash",
            contents=contents_payload,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=LLMEvaluationResult,
                temperature=0.3
            )
        )

        llm_result = LLMEvaluationResult.model_validate_json(response.text)

        # Save evaluation results to DB
        db_eval = InterviewEvaluation(
            interview_id=interview.interview_id,
            overall_score=llm_result.overall_score,
            general_feedback=llm_result.general_feedback,
            evaluations_json=[e.model_dump() for e in llm_result.evaluations]
        )
        db.add(db_eval)
        db.commit()
        db.refresh(db_eval)

        # Map response
        response_evals = []
        for e in llm_result.evaluations:
            response_evals.append(QuestionEvaluationItem(
                question_id=e.question_id,
                score=e.score,
                rationale=e.rationale,
                model_answer=e.model_answer,
                improved_phrasing=e.improved_phrasing
            ))

        return EvaluateInterviewResponse(
            overall_score=db_eval.overall_score,
            general_feedback=db_eval.general_feedback,
            evaluations=response_evals
        )

    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to evaluate defense answers: {str(e)}")


@router.post("/parse-file")
async def parse_uploaded_file(file: UploadFile = File(...)):
    # Fallback endpoint if needed for UI details
    filename = file.filename.lower()
    valid_extensions = (".pdf", ".txt", ".docx", ".doc", ".pptx", ".ppt")
    if not filename.endswith(valid_extensions):
        raise HTTPException(status_code=400, detail="Only PDF, TXT, Word (.docx) and PowerPoint (.pptx) files are supported")
        
    try:
        text = await extract_text_from_upload_file(file)
        return {"text": text.strip()}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse file: {str(e)}")
